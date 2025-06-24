from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Slide content
slides_content = [
    # (Title, Bullets/List)
    ("Transforming Construction Billing",
     ["BoQ Management App | Investor Presentation | [Month, Year]",
      "Presented by: [Your Name, Contact]"]),
    
    ("Why Construction Billing is Broken",
     ["80% of construction firms still use Excel for BoQ and billing",
      "Manual processes → delays, errors, payment disputes",
      "Milestone billing rarely automated; cash flow gaps",
      "Regulations (GST, VAT, BIM) demand digital compliance"]),

    ("Our Platform: Modern, Digital, Automated",
     ["Cloud-based BoQ builder & Excel upload",
      "Milestone-based payment breakups for every item",
      "Work progress tracking, instant bill generation",
      "Audit-ready, integration-ready, easy to use"]),

    ("A Global, Underserved Market",
     ["Global construction IT spend: $180B+/year (Statista)",
      "SaaS construction software CAGR: 9.3% (Markets&Markets)",
      "TAM: $2.5–4B globally; SAM: $500M in 6 core markets",
      "[Insert map/chart: Market size by country]"]),

    ("Why We Win (Competitive Landscape)",
     ["Procore, Autodesk: Expensive, complex, US-centric",
      "Excel/manual: Prone to errors, non-collaborative",
      "Our App: Simple, affordable, focused on billing automation",
      "[Insert 2x2 Matrix: Price vs Usability]"]),

    ("How It Works (Product Demo Flow)",
     ["Register/Login (secure, multi-role)",
      "Create/upload BoQ, set payment breakups",
      "Track progress, upload work completion, generate bills",
      "Dashboard analytics, export, integration",
      "[Insert workflow screenshots or diagram]"]),

    ("How We Make Money (Business Model)",
     ["SaaS subscription (per project, per org)",
      "Onboarding & professional services",
      "Integration add-ons (ERP, accounting, GST/VAT modules)",
      "[Insert pie chart: Revenue by stream]"]),

    ("Countrywise Go-To-Market",
     ["UK/EU (pilot): LinkedIn, RICS, trade shows, PM consultants",
      "India: EPCs, infra, CA/accountant partners, digital onboarding",
      "USA: SME contractors, channel partners",
      "UAE/Australia: Compliance-driven, via local agents",
      "[Insert map with entry arrows/flowchart]"]),

    ("Execution Plan & Timeline",
     ["Q1: Design & Foundation (UI/UX, tech setup)",
      "Q2: MVP Build (core features, billing, dashboard)",
      "Q3: Pilots & Feedback (UK/India focus)",
      "Q4: Integrations, Launch, Global Expansion",
      "[Insert Gantt/timeline chart]"]),

    ("Financial Projections",
     ["Year 1: 50 clients, £60K–£100K ARR",
      "Year 2: 250+ clients, £300K+ ARR",
      "[Insert ARR growth line chart; client bar chart]"]),

    ("Execution-Driven Leadership (Team)",
     ["[Your Name] — Founder/Product Lead (Bio, LinkedIn)",
      "[CTO/Backend Lead]", "[Frontend Lead]",
      "[QA/Support]", "[Advisors: industry, legal, finance]",
      "[Insert team headshots/bios]"]),

    ("Built for Resilience (Risks & Mitigation)",
     ["Adoption: Free onboarding, super-simple UI",
      "Security: Pen-testing, GDPR, encrypted data",
      "Regulatory: Modular compliance, local partners",
      "[Insert risk/mitigation table]"]),

    ("Join Us in Transforming Construction (The Ask)",
     ["Seeking £150,000 pre-seed for 12-month runway",
      "Funds: 60% dev, 20% GTM, 10% infra, 10% ops",
      "Equity: 10–15% (negotiable)",
      "Next steps: Book a demo, Q&A",
      "[Insert fund allocation pie chart]"]),

    ("Let’s Build the Future Together",
     ["Contact: [Email, LinkedIn, Demo Link]",
      "Thank You!"])
]

def add_slide(prs, title, lines):
    slide_layout = prs.slide_layouts[1]  # Title & Content
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    tf = content_placeholder.text_frame
    for i, line in enumerate(lines):
        if i == 0:
            tf.text = line
        else:
            tf.add_paragraph().text = line

prs = Presentation()
for title, lines in slides_content:
    add_slide(prs, title, lines)

prs.save("BoQ_Investor_Deck.pptx")
print("Presentation saved as BoQ_Investor_Deck.pptx")
